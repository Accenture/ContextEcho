"""ContextEcho framework figure — v3, modeled after Firefly mockup.

Layout (mirrors the Firefly reference):
  Top banner (full width):
    Our Vision (left) | 4 principles inline (center) | Our Goal (right)
  Body row:
    Left sidebar: Evaluation Coverage (orgs + sessions) | 4 stimuli cards
  Primitive band (full width below sidebar+cards):
    Title + 3 donor sessions (left) | equations (center) | "12 Snapshot Positions" (left of equations)
  Bottom strip (NEW, replaces Q1-Q5):
    3 grouped metric columns | Impact & Vision panel
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


def hexc(s):
    s = s.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


# ---------- palette (matches Firefly's clean look) ----------
C_PROBES = hexc("16a34a")
C_STRESSORS = hexc("7c3aed")
C_SWEBENCH = hexc("ea580c")
C_ANCHOR = hexc("dc2626")

C_CANVAS = hexc("ffffff")
C_PRIM_BG = hexc("fffbeb")
C_PRIM_EC = hexc("f59e0b")
C_FRAME_BG = hexc("f8fafc")
C_FRAME_EC = hexc("e2e8f0")
C_SIDEBAR_BG = hexc("f9fafb")
C_SIDEBAR_EC = hexc("e5e7eb")
C_TOP_BG = hexc("eff6ff")
C_TOP_EC = hexc("bfdbfe")

C_HEADER = hexc("0f172a")
C_HEADING_BLUE = hexc("1e3a8a")
C_BODY = hexc("334155")
C_MUTED = hexc("64748b")
C_FAINT = hexc("94a3b8")
C_WHITE = hexc("ffffff")


# ---------- helpers ----------
def add_rect(slide, x, y, w, h, fill, line=None, lw=0.5, rounded=False, corner=0.05):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    if rounded:
        sh.adjustments[0] = corner
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def add_text(slide, x, y, w, h, text, *, size=11, bold=False, italic=False,
             color=None, align="left", anchor="top", font="Calibri", spacing=None):
    if color is None:
        color = C_BODY
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
        if spacing is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(int(spacing * 100)))
    return tb


def add_line(slide, x1, y1, x2, y2, color, lw=1.0):
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(lw)
    return ln


def icon_circle(slide, cx, cy, r, color):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - r), Inches(cy - r), Inches(r * 2), Inches(r * 2),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


# ---------- glyphs (white-on-circle) ----------
def glyph_target(slide, cx, cy, r, color):
    icon_circle(slide, cx, cy, r, color)
    for ratio, fill in [(0.55, C_WHITE), (0.35, color), (0.15, C_WHITE)]:
        sh = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(cx - r * ratio), Inches(cy - r * ratio),
            Inches(r * ratio * 2), Inches(r * ratio * 2),
        )
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        sh.line.fill.background()
        sh.shadow.inherit = False


def glyph_balance(slide, cx, cy, r, color):
    icon_circle(slide, cx, cy, r, color)
    add_line(slide, cx - r * 0.55, cy - r * 0.10, cx + r * 0.55, cy - r * 0.10, C_WHITE, lw=2.5)
    add_line(slide, cx, cy - r * 0.10, cx, cy + r * 0.40, C_WHITE, lw=2.5)
    add_line(slide, cx - r * 0.30, cy + r * 0.40, cx + r * 0.30, cy + r * 0.40, C_WHITE, lw=2.5)


def glyph_chart(slide, cx, cy, r, color):
    icon_circle(slide, cx, cy, r, color)
    bar_w = r * 0.18
    heights = [0.28, 0.45, 0.62]
    base_y = cy + r * 0.35
    starts_x = [cx - r * 0.45, cx - r * 0.10, cx + r * 0.25]
    for h, sx in zip(heights, starts_x):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(sx), Inches(base_y - r * h),
            Inches(bar_w), Inches(r * h),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = C_WHITE
        bar.line.fill.background()
        bar.shadow.inherit = False


def glyph_lock(slide, cx, cy, r, color):
    icon_circle(slide, cx, cy, r, color)
    shackle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - r * 0.30), Inches(cy - r * 0.50),
        Inches(r * 0.60), Inches(r * 0.55),
    )
    shackle.fill.background()
    shackle.line.color.rgb = C_WHITE
    shackle.line.width = Pt(2.5)
    shackle.shadow.inherit = False
    body = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(cx - r * 0.40), Inches(cy - r * 0.15),
        Inches(r * 0.80), Inches(r * 0.55),
    )
    body.adjustments[0] = 0.15
    body.fill.solid()
    body.fill.fore_color.rgb = C_WHITE
    body.line.fill.background()
    body.shadow.inherit = False


def glyph_eye(slide, cx, cy, r, color):
    icon_circle(slide, cx, cy, r, color)
    eye = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - r * 0.55), Inches(cy - r * 0.30),
        Inches(r * 1.10), Inches(r * 0.60),
    )
    eye.fill.solid()
    eye.fill.fore_color.rgb = C_WHITE
    eye.line.fill.background()
    eye.shadow.inherit = False
    pup = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - r * 0.20), Inches(cy - r * 0.20),
        Inches(r * 0.40), Inches(r * 0.40),
    )
    pup.fill.solid()
    pup.fill.fore_color.rgb = color
    pup.line.fill.background()
    pup.shadow.inherit = False


def glyph_terminal(slide, cx, cy, r, color):
    icon_circle(slide, cx, cy, r, color)
    add_line(slide, cx - r * 0.40, cy - r * 0.30, cx - r * 0.05, cy, C_WHITE, lw=3.0)
    add_line(slide, cx - r * 0.05, cy, cx - r * 0.40, cy + r * 0.30, C_WHITE, lw=3.0)
    add_line(slide, cx + r * 0.05, cy + r * 0.30, cx + r * 0.40, cy + r * 0.30, C_WHITE, lw=3.0)


def glyph_anchor(slide, cx, cy, r, color):
    """Anchor symbol for A-anchor mitigation."""
    icon_circle(slide, cx, cy, r, color)
    # ring at top
    ring = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - r * 0.18), Inches(cy - r * 0.55),
        Inches(r * 0.36), Inches(r * 0.30),
    )
    ring.fill.background()
    ring.line.color.rgb = C_WHITE
    ring.line.width = Pt(2.4)
    ring.shadow.inherit = False
    # vertical post
    add_line(slide, cx, cy - r * 0.25, cx, cy + r * 0.45, C_WHITE, lw=3.0)
    # crossbar
    add_line(slide, cx - r * 0.35, cy - r * 0.10, cx + r * 0.35, cy - r * 0.10, C_WHITE, lw=2.5)
    # hooks at bottom
    add_line(slide, cx - r * 0.40, cy + r * 0.45, cx, cy + r * 0.55, C_WHITE, lw=2.5)
    add_line(slide, cx + r * 0.40, cy + r * 0.45, cx, cy + r * 0.55, C_WHITE, lw=2.5)


def glyph_chevron_right(slide, cx, cy, r, color):
    """For Our Goal: target + arrow."""
    icon_circle(slide, cx, cy, r, color)
    # inner target ring
    o = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - r * 0.45), Inches(cy - r * 0.45),
        Inches(r * 0.90), Inches(r * 0.90),
    )
    o.fill.background()
    o.line.color.rgb = C_WHITE
    o.line.width = Pt(1.8)
    o.shadow.inherit = False
    # bullseye
    b = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - r * 0.15), Inches(cy - r * 0.15),
        Inches(r * 0.30), Inches(r * 0.30),
    )
    b.fill.solid()
    b.fill.fore_color.rgb = C_WHITE
    b.line.fill.background()
    b.shadow.inherit = False


# ---------- builder ----------
def build(out_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    pad = 0.12
    W = 13.333
    H = 7.5

    # ============== TOP BANNER (mirrors Firefly) ==============
    top_h = 0.90
    top_y = pad
    add_rect(slide, pad, top_y, W - 2 * pad, top_h, C_TOP_BG, C_TOP_EC, lw=1.0, rounded=True, corner=0.02)

    # Our Vision (far left, ~1.6in wide)
    vis_w = 1.65
    vx = pad + 0.12
    glyph_target(slide, vx + 0.20, top_y + 0.30, 0.14, C_HEADING_BLUE)
    add_text(slide, vx + 0.50, top_y + 0.08, vis_w - 0.50, 0.32,
             "OUR VISION", size=11, bold=True, color=C_HEADING_BLUE, anchor="middle", spacing=1.5)
    add_text(slide, vx, top_y + 0.45, vis_w + 0.20, 0.42,
             "A reusable benchmark for persona drift\nin long agentic-coding sessions.",
             size=8.5, color=C_MUTED, anchor="top")
    # divider
    add_line(slide, vx + vis_w + 0.10, top_y + 0.10, vx + vis_w + 0.10, top_y + top_h - 0.10, C_FAINT, lw=0.8)

    # Our Goal (far right)
    goal_w = 1.85
    goal_x = W - pad - goal_w - 0.12
    add_rect(slide, goal_x, top_y + 0.10, goal_w, top_h - 0.20, C_WHITE, C_HEADING_BLUE, lw=1.2, rounded=True, corner=0.04)
    glyph_chevron_right(slide, goal_x + 0.22, top_y + top_h / 2, 0.14, C_HEADING_BLUE)
    # arrow connecting to box from outside
    add_line(slide, goal_x - 0.10, top_y + top_h / 2, goal_x - 0.02, top_y + top_h / 2, C_HEADING_BLUE, lw=2.0)
    add_text(slide, goal_x + 0.50, top_y + 0.12, goal_w - 0.55, 0.32,
             "OUR GOAL", size=11, bold=True, color=C_HEADING_BLUE, anchor="middle", spacing=1.5)
    add_text(slide, goal_x + 0.50, top_y + 0.42, goal_w - 0.55, 0.50,
             "Audit whether the persona a model\nships with is the persona users\nencounter at session end.",
             size=8, color=C_MUTED, anchor="top")

    # Center: MEASUREMENT STIMULI title + 4 principles
    center_x_start = vx + vis_w + 0.25
    center_x_end = goal_x - 0.25
    center_w = center_x_end - center_x_start

    add_text(slide, center_x_start, top_y + 0.04, center_w, 0.25,
             "MEASUREMENT STIMULI (STIMULI LAYER)",
             size=12, bold=True, color=C_HEADING_BLUE, align="center")

    principles = [
        (glyph_target, "Snapshot-then-probe", "Fork the session at turn t,\nprobe without perturbing"),
        (glyph_balance, "Length-matched filler", "Lorem-ipsum control strips\nthe family signal"),
        (glyph_chart, "Dual measurement", "Judge-scored probes plus\njudge-free regex compliance"),
        (glyph_lock, "Pre-registered analysis", "SHA-256 hashed plans\nlocked before data lands"),
    ]
    n_p = 4
    p_w = center_w / n_p
    p_y = top_y + 0.36
    for i, (g, title, sub) in enumerate(principles):
        px = center_x_start + i * p_w
        g(slide, px + 0.16, p_y + 0.18, 0.12, C_HEADING_BLUE)
        add_text(slide, px + 0.38, p_y + 0.02, p_w - 0.42, 0.22,
                 title, size=10, bold=True, color=C_HEADING_BLUE)
        add_text(slide, px + 0.38, p_y + 0.22, p_w - 0.42, 0.32,
                 sub, size=8.0, color=C_MUTED)

    # ============== LEFT SIDEBAR ==============
    side_w = 1.85
    side_top = top_y + top_h + 0.18
    prim_h = 1.20
    prim_y_target = H - pad - 1.10 - 0.20 - prim_h  # primitive sits above bottom strip
    side_bot = prim_y_target  # sidebar runs down to primitive band
    side_h = side_bot - side_top
    add_rect(slide, pad, side_top, side_w, side_h, C_SIDEBAR_BG, C_SIDEBAR_EC, lw=0.8, rounded=True, corner=0.02)

    # Heading
    add_text(slide, pad + 0.10, side_top + 0.12, side_w - 0.20, 0.28,
             "EVALUATION COVERAGE", size=10, bold=True, color=C_HEADING_BLUE, align="center", spacing=1.5)
    add_text(slide, pad + 0.10, side_top + 0.40, side_w - 0.20, 0.34,
             "23 frontier targets\nfrom 10 organizations",
             size=8, italic=True, color=C_MUTED, align="center")

    # Org list with logos as colored dots
    org_y_start = side_top + 0.92
    orgs = [
        ("Anthropic", "Haiku, Sonnet 4.5/4.6,\nOpus 4.1"),
        ("OpenAI", "GPT-4.1, 4o, 5, 5-mini"),
        ("Google", "Gemini 2.5 Pro, Flash"),
        ("DeepSeek", "V3"),
        ("Mistral", "Small, Medium, Large"),
        ("Cohere", "Command R7B, Command A"),
        ("NVIDIA", "Nemotron Nano, 49B, 120B"),
        ("Alibaba", "Qwen3 235B, Next 80B"),
        ("Meta", "Llama 3.3 70B"),
        ("Moonshot", "Kimi K2.6"),
    ]
    row_h = (side_h - 1.05) / len(orgs)
    for i, (org, models) in enumerate(orgs):
        ry = org_y_start + i * row_h
        # logo placeholder dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(pad + 0.14), Inches(ry + 0.03),
            Inches(0.20), Inches(0.20),
        )
        dot.adjustments[0] = 0.30
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_HEADING_BLUE
        dot.line.fill.background()
        dot.shadow.inherit = False
        add_text(slide, pad + 0.40, ry, 0.85, 0.22,
                 org, size=8.5, bold=True, color=C_HEADER, anchor="top")
        add_text(slide, pad + 0.40, ry + 0.17, side_w - 0.55, 0.30,
                 models.split("\n")[0], size=7.5, color=C_MUTED)

    # ============== "Probe Invocation" arrows between sidebar and cards ==============
    arrow_x1 = pad + side_w + 0.01
    arrow_x2 = pad + side_w + 0.45
    arrow_y_top = side_top + 1.0
    arrow_y_bot = side_top + 1.6

    # right arrow (probe invocation)
    add_line(slide, arrow_x1, arrow_y_top, arrow_x2, arrow_y_top, C_MUTED, lw=1.5)
    tri1 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        Inches(arrow_x2 - 0.02), Inches(arrow_y_top - 0.08),
        Inches(0.10), Inches(0.16),
    )
    tri1.rotation = 90
    tri1.fill.solid()
    tri1.fill.fore_color.rgb = C_MUTED
    tri1.line.fill.background()
    tri1.shadow.inherit = False
    add_text(slide, arrow_x1, arrow_y_top - 0.30, 0.50, 0.20,
             "Probe\nInvocation", size=7, italic=True, color=C_MUTED, align="center", anchor="middle")

    # left arrow (structured response)
    add_line(slide, arrow_x2, arrow_y_bot, arrow_x1, arrow_y_bot, C_MUTED, lw=1.5)
    tri2 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        Inches(arrow_x1 - 0.05), Inches(arrow_y_bot - 0.08),
        Inches(0.10), Inches(0.16),
    )
    tri2.rotation = 270
    tri2.fill.solid()
    tri2.fill.fore_color.rgb = C_MUTED
    tri2.line.fill.background()
    tri2.shadow.inherit = False
    add_text(slide, arrow_x1, arrow_y_bot + 0.10, 0.50, 0.20,
             "Structured\nResponse", size=7, italic=True, color=C_MUTED, align="center", anchor="middle")

    # ============== 4 STIMULI CARDS ==============
    main_x = pad + side_w + 0.55
    main_w = W - main_x - pad
    stim_y = side_top
    stim_h = side_h  # equals sidebar height

    # row label above cards
    add_text(slide, main_x, stim_y - 0.20, main_w, 0.18,
             "— One Primitive, Many Stimuli",
             size=10, italic=True, color=C_MUTED, align="center")

    cards = [
        dict(color=C_PROBES, num="1", title="IDENTITY PROBES", glyph=glyph_eye,
             tagline="judge-scored, granular",
             body="25-probe off-task battery scored\non a 4-pt assistant-register rubric.",
             kvs=[("CATEGORIES:", "Identity 4 · Experience 8\nPreference 4 · Relational 4\nCoding-Self 5 (primary)"),
                  ("PARAPHRASES:", "n=10 per cell")]),
        dict(color=C_STRESSORS, num="2", title="FORMAT STRESSORS", glyph=glyph_balance,
             tagline="judge-free, deployment-relevant",
             body="Four format-constraint instructions\nscored by deterministic regex.",
             kvs=[("SCOPE:", "S1 byte-exact 1-word\nS2 no-preamble (primary)\nS3 1-sentence · S4 byte-exact JSON"),
                  ("METRIC:", "% compliance, length ratio")]),
        dict(color=C_SWEBENCH, num="3", title="SWE-BENCH CONTINUATION", glyph=glyph_terminal,
             tagline="tool-using, paired",
             body="Agent shown Claude- vs.\nGPT-flavored prefix → next tool call.",
             kvs=[("COVERAGE:", "25 cutpoints × 3 targets:\nSonnet 4.6, Mistral Small,\nKimi K2.6"),
                  ("METRIC:", "paired Δ argument fidelity")]),
        dict(color=C_ANCHOR, num="4", title="A-ANCHOR MITIGATION", glyph=glyph_anchor,
             tagline="single-shot, deployment-grade",
             body="~80-token user-turn block\ninserted between prefix and probe.",
             kvs=[("RECIPES:", "V0 identity (30t) · V2 format demo\nV0+V2 combined · 200t large variant"),
                  ("PERSISTENCE:", "≥20 unanchored turns immunized")]),
    ]

    n_cards = 4
    card_gap = 0.14
    card_w = (main_w - card_gap * (n_cards - 1)) / n_cards

    for i, card in enumerate(cards):
        x = main_x + i * (card_w + card_gap)
        add_rect(slide, x, stim_y, card_w, stim_h, C_CANVAS,
                 line=card["color"], lw=1.8, rounded=True, corner=0.025)

        inner_x = x + 0.18
        inner_w = card_w - 0.36

        # numbered title
        add_text(slide, inner_x, stim_y + 0.12, inner_w, 0.28,
                 f"{card['num']}. {card['title']}", size=12, bold=True, color=card["color"])

        # tagline
        add_text(slide, inner_x, stim_y + 0.40, inner_w, 0.20,
                 card["tagline"], size=8.5, italic=True, color=C_MUTED)

        # glyph centered
        card["glyph"](slide, x + card_w / 2, stim_y + 0.95, 0.24, card["color"])

        # body description
        add_text(slide, inner_x, stim_y + 1.40, inner_w, 0.60,
                 card["body"], size=8.8, color=C_BODY, align="center")

        # divider
        add_line(slide, inner_x + 0.20, stim_y + 2.05, x + card_w - 0.20 - 0.20, stim_y + 2.05, card["color"], lw=0.6)

        # key-value blocks
        kv_y = stim_y + 2.15
        for k, v in card["kvs"]:
            add_text(slide, inner_x, kv_y, inner_w, 0.20,
                     k, size=8, bold=True, color=card["color"], spacing=1)
            add_text(slide, inner_x, kv_y + 0.20, inner_w, 0.65,
                     v, size=7.5, color=C_BODY)
            kv_y += 0.95

    # ============== PRIMITIVE BAND — 3-column internal layout ==============
    prim_y = side_bot + 0.18
    add_rect(slide, pad, prim_y, W - 2 * pad, prim_h, C_PRIM_BG,
             line=C_PRIM_EC, lw=1.2, rounded=True, corner=0.02)

    # Title row (spans full width)
    add_text(slide, pad + 0.20, prim_y + 0.08, 4.0, 0.26,
             "Snapshot-then-Probe Primitive",
             size=13, bold=True, color=C_HEADER)
    add_text(slide, pad + 0.20 + 3.40, prim_y + 0.12, W - 2 * pad - 3.80, 0.22,
             "— shared infrastructure: every stimulus flows through this primitive",
             size=9, italic=True, color=C_MUTED)

    # 3-column layout below title
    col_top = prim_y + 0.40
    col_h = prim_h - 0.66

    # Column boundaries
    left_x = pad + 0.20
    left_w = 3.05
    right_w = 3.05
    right_x = W - pad - 0.20 - right_w
    center_x = left_x + left_w + 0.20
    center_w = right_x - center_x - 0.20

    # ----- Left column: 3 Donor Sessions -----
    add_text(slide, left_x, col_top, left_w, 0.20,
             "3 Donor Sessions",
             size=9.5, bold=True, color=C_HEADING_BLUE)
    sessions = [
        ("S1:", "9,643 turns", "agentic coding"),
        ("S2:", "3,746 turns", "manuscript writing"),
        ("S3:", "4,918 turns", "non-coding docs"),
    ]
    sess_y = col_top + 0.22
    for lab, turns, kind in sessions:
        add_text(slide, left_x, sess_y, 0.35, 0.18,
                 lab, size=8.5, bold=True, color=C_HEADER)
        add_text(slide, left_x + 0.35, sess_y, 0.95, 0.18,
                 turns, size=8.5, color=C_BODY)
        add_text(slide, left_x + 1.32, sess_y, left_w - 1.32, 0.18,
                 kind, size=7.8, italic=True, color=C_MUTED)
        sess_y += 0.18

    # ----- Center column: equations -----
    add_text(slide, center_x, col_top + 0.02, center_w, 0.28,
             "r⁽ᶜˡᵃᵘᵈᵉ⁾ₜ,ₚ,ᵢ = ℳ(c⁽ᶜˡᵃᵘᵈᵉ⁾₁:ₜ ⊕ FRAME ⊕ pᵢ)       "
             "r⁽ᶠⁱˡˡᵉʳ⁾ₜ,ₚ,ᵢ = ℳ(c⁽ᶠⁱˡˡᵉʳ⁾₁:ₜ ⊕ FRAME ⊕ pᵢ)",
             size=10.5, italic=True, color=C_HEADER, font="Cambria Math")
    add_text(slide, center_x, col_top + 0.34, center_w, 0.28,
             "Δ(t, ℳ) = (1/|P|) Σₚ∈ₚ [ J̄(r⁽ᶠⁱˡˡᵉʳ⁾ₜ,ₚ) − J̄(r⁽ᶜˡᵃᵘᵈᵉ⁾ₜ,ₚ) ]    "
             "Δ > 0 ⇒ persona drift",
             size=10.5, italic=True, color=C_HEADER, font="Cambria Math")

    # ----- Right column: 12 Snapshot Positions with single-row timeline -----
    add_text(slide, right_x, col_top, right_w, 0.20,
             "12 Snapshot Positions",
             size=9.5, bold=True, color=C_HEADING_BLUE)

    # Single-row timeline: 12 dots, vertical divider between pre/post
    n_dots = 12
    dot_r = 0.030
    track_left = right_x + 0.05
    track_right = right_x + right_w - 0.05
    track_w = track_right - track_left
    dot_step = track_w / (n_dots - 1)
    track_y = col_top + 0.38

    # base line
    add_line(slide, track_left, track_y, track_right, track_y, C_HEADING_BLUE, lw=0.8)

    # vertical divider between dot 5 (P5, last pre) and dot 6 (C1, first post)
    # Position the divider midway between the two
    divider_x = track_left + 5.5 * dot_step
    add_line(slide, divider_x, track_y - 0.10, divider_x, track_y + 0.10, C_PRIM_EC, lw=1.2)

    # 12 dots (first 6 blue = pre, last 6 amber = post)
    pre_color = C_HEADING_BLUE
    post_color = C_PRIM_EC
    for i in range(n_dots):
        cx = track_left + i * dot_step
        color = pre_color if i < 6 else post_color
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(cx - dot_r), Inches(track_y - dot_r),
            Inches(dot_r * 2), Inches(dot_r * 2),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        dot.shadow.inherit = False

    # group labels above the timeline (pre / post)
    add_text(slide, track_left - 0.05, col_top + 0.22, 2.5 * dot_step + 0.10, 0.14,
             "pre-compaction", size=6.8, italic=True, color=C_HEADING_BLUE, align="center")
    add_text(slide, track_left + 6 * dot_step - 0.05, col_top + 0.22, 5 * dot_step + 0.10, 0.14,
             "post-compaction", size=6.8, italic=True, color=C_PRIM_EC, align="center")

    # P0..P5 and C1..C6 labels below the dots, with very compact font
    labels = ["P0", "P1", "P2", "P3", "P4", "P5", "C1", "C2", "C3", "C4", "C5", "C6"]
    for i, lab in enumerate(labels):
        cx = track_left + i * dot_step
        add_text(slide, cx - 0.10, track_y + 0.05, 0.20, 0.14,
                 lab, size=6.0, color=C_MUTED, align="center")

    # ----- footer (centered, across whole band) -----
    add_text(slide, pad + 0.20, prim_y + prim_h - 0.26, W - 2 * pad - 0.40, 0.20,
             "Fork is discarded after sampling — the main session is never perturbed.",
             size=8.5, italic=True, color=C_MUTED, align="center")

    # ============== BOTTOM STRIP — grouped metrics + Impact panel ==============
    bot_h = H - pad - (prim_y + prim_h + 0.18)
    bot_y = prim_y + prim_h + 0.18

    metric_w = (W - 2 * pad) * 0.78
    impact_x = pad + metric_w + 0.18
    impact_w = W - pad - impact_x

    # Metrics container
    add_rect(slide, pad, bot_y, metric_w, bot_h, C_FRAME_BG,
             line=C_FRAME_EC, lw=0.8, rounded=True, corner=0.02)
    add_text(slide, pad, bot_y + 0.06, metric_w, 0.22,
             "Evaluation Outcomes",
             size=12, bold=True, color=C_HEADER, align="center")
    add_text(slide, pad, bot_y + 0.26, metric_w, 0.18,
             "drift detection across complementary surfaces",
             size=8.5, italic=True, color=C_MUTED, align="center")

    groups = [
        dict(title="Probe Surface", sub="judge-scored", color=C_PROBES,
             metrics=[("Drift gap Δ", "0–3 rubric"),
                      ("Per-category Δ", "5 categories"),
                      ("Behavioral fingerprint", "PCA, judge-blind")]),
        dict(title="Stressor Surface", sub="judge-free", color=C_STRESSORS,
             metrics=[("Compliance rate", "S₂ regex"),
                      ("Length ratio", "vs filler"),
                      ("Compaction trajectory", "pre/post Δ")]),
        dict(title="Downstream", sub="deployment cost", color=C_SWEBENCH,
             metrics=[("Argument fidelity", "paired Δ"),
                      ("Contract break", "Δpp compliance"),
                      ("Cost direction", "mode-dependent")]),
    ]
    g_gap = 0.18
    g_inset = pad + 0.20
    g_total_w = metric_w - 0.40
    g_w = (g_total_w - g_gap * (len(groups) - 1)) / len(groups)
    g_y = bot_y + 0.50
    g_h = bot_h - 0.60

    for i, g in enumerate(groups):
        gx = g_inset + i * (g_w + g_gap)
        # vertical accent
        add_rect(slide, gx, g_y, 0.06, g_h, g["color"], line=None)
        # title + sub
        add_text(slide, gx + 0.15, g_y, g_w - 0.15, 0.22,
                 g["title"], size=10.5, bold=True, color=C_HEADER)
        add_text(slide, gx + 0.15, g_y + 0.20, g_w - 0.15, 0.18,
                 g["sub"], size=8.5, italic=True, color=C_MUTED)
        # metrics
        m_y = g_y + 0.46
        for name, unit in g["metrics"]:
            sq = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(gx + 0.16), Inches(m_y + 0.06),
                Inches(0.08), Inches(0.08),
            )
            sq.fill.solid()
            sq.fill.fore_color.rgb = g["color"]
            sq.line.fill.background()
            sq.shadow.inherit = False
            add_text(slide, gx + 0.30, m_y, g_w - 0.30, 0.18,
                     name, size=9, bold=True, color=C_BODY)
            add_text(slide, gx + 0.30, m_y + 0.16, g_w - 0.30, 0.18,
                     unit, size=7.5, italic=True, color=C_MUTED)
            m_y += 0.36

    # Impact & Vision panel
    add_rect(slide, impact_x, bot_y, impact_w, bot_h, C_CANVAS,
             line=C_HEADING_BLUE, lw=1.2, rounded=True, corner=0.04)
    # chevron icon
    chev = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON,
        Inches(impact_x + 0.18), Inches(bot_y + 0.15),
        Inches(0.32), Inches(0.24),
    )
    chev.fill.solid()
    chev.fill.fore_color.rgb = C_HEADING_BLUE
    chev.line.fill.background()
    chev.shadow.inherit = False
    add_text(slide, impact_x + 0.56, bot_y + 0.15, impact_w - 0.60, 0.25,
             "IMPACT & VISION", size=10, bold=True, color=C_HEADING_BLUE, spacing=1.5)
    add_text(slide, impact_x + 0.18, bot_y + 0.48, impact_w - 0.36, bot_h - 0.55,
             "Open-source benchmark that audits whether "
             "the shipped persona is the persona users "
             "encounter at session end.",
             size=8, color=C_BODY, anchor="top")

    prs.save(out_path)
    print(f"wrote {out_path}")


if __name__ == "__main__":
    here = Path(__file__).parent
    out = here.parent / "paper" / "figures" / "fig_framework.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    build(out)
